from pptx import Presentation
from pptx.util import Inches

# Create a Presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Food Safety and Quality"
subtitle.text = "An Overview of Practices and Metrics"

# Slide 2: Inspection Results
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Inspection Results"

content = (
    "The following table presents the inspection results for various food products over the past week.\n\n"
    "Date          Product        Inspection Type        Result        Inspector\n"
    "2024-07-01    Apple          Regular                Pass          John Doe\n"
    "2024-07-02    Banana         Follow-up              Pass          Jane Smith\n"
    "2024-07-03    Carrot         Regular                Fail          Emily Johnson\n"
    "2024-07-04    Spinach        Regular                Pass          Michael Brown\n"
    "2024-07-05    Strawberry     Complaint              Fail          Chris Davis\n"
)
textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4))
text_frame = textbox.text_frame
text_frame.text = content

# Slide 3: Quality Metrics
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Quality Metrics"

content = (
    "The following metrics provide an overview of the quality of different food products inspected.\n\n"
    "1. Color: Apples are assessed for their color, with a standard range of 7-10.\n"
    "2. Firmness: Bananas are assessed for firmness, with a standard range of 3-6.\n"
    "3. Moisture Level: Carrots are assessed for moisture level, with a standard range of 10-15%.\n"
    "4. Pesticide Level: Spinach is assessed for pesticide level, with a standard of <=0.02%.\n"
    "5. Brix Level: Strawberries are assessed for Brix level, with a standard range of 8-12.\n"
)
textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4))
text_frame = textbox.text_frame
text_frame.text = content

# Slide 4: Best Practices
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Best Practices for Food Safety"

content = (
    "1. Hand Washing: Ensure all personnel wash their hands thoroughly before handling food.\n"
    "2. Proper Storage: Store food at the correct temperatures to prevent spoilage and contamination.\n"
    "3. Cross-contamination Prevention: Use separate equipment and surfaces for raw and cooked foods.\n"
    "4. Regular Inspections: Conduct regular inspections to ensure compliance with food safety standards.\n"
    "5. Training: Provide ongoing training for staff on food safety practices and protocols.\n"
)
textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4))
text_frame = textbox.text_frame
text_frame.text = content

# Save the presentation
prs.save('food_safety_presentation.pptx')
